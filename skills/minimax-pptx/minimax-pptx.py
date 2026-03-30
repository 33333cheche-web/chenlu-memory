#!/usr/bin/env python3
"""
MiniMax-pptx Skill - PPT 创建和编辑
支持：创建演示文稿、添加幻灯片、设置样式
"""

import argparse
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_ppt(title, slides_data, output_path):
    """创建PPT文件"""
    prs = Presentation()
    
    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 标题页
    title_slide_layout = prs.slide_layouts[0]  # 标题页布局
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "上海沃橙信息技术有限公司"
    
    # 内容页
    for slide_data in slides_data:
        content_slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = slide_data.get('title', '')
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = slide_data.get('content', '')
    
    prs.save(output_path)
    return output_path

def main():
    parser = argparse.ArgumentParser(description='MiniMax-pptx: PPT创建和编辑')
    parser.add_argument('--title', '-t', type=str, required=True, help='PPT标题')
    parser.add_argument('--slides', '-s', type=str, help='JSON格式幻灯片数据')
    parser.add_argument('--output', '-o', type=str, required=True, help='输出路径')
    
    args = parser.parse_args()
    
    import json
    slides = json.loads(args.slides) if args.slides else []
    output_path = create_ppt(args.title, slides, args.output)
    print(f"✅ PPT已创建: {output_path}")

if __name__ == "__main__":
    main()
